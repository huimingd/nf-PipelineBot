from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── palette ──────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1E, 0x1E, 0x2E)   # slide background
ACCENT     = RGBColor(0x89, 0xB4, 0xFA)   # Catppuccin blue
GREEN      = RGBColor(0xA6, 0xE3, 0xA1)   # Catppuccin green
YELLOW     = RGBColor(0xF9, 0xE2, 0xAF)   # Catppuccin yellow
MAUVE      = RGBColor(0xCB, 0xA6, 0xF7)   # Catppuccin mauve
RED        = RGBColor(0xF3, 0x8B, 0xA8)   # Catppuccin red
TEXT_MAIN  = RGBColor(0xCD, 0xD6, 0xF4)   # light text
TEXT_SUB   = RGBColor(0xA6, 0xAD, 0xC8)   # muted text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BOX_BG     = RGBColor(0x31, 0x32, 0x44)   # card background

# ── slide size: widescreen 16:9 ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── helpers ──────────────────────────────────────────────────────────────────
def new_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return sl


def add_text(sl, text, x, y, w, h, size=18, bold=False, color=TEXT_MAIN,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_box(sl, x, y, w, h, fill=BOX_BG, line_color=ACCENT, line_width_pt=1):
    shape = sl.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width_pt)
    return shape


def slide_header(sl, title, subtitle=None):
    """Top bar with title + optional subtitle."""
    bar = sl.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x18, 0x18, 0x2E)
    bar.line.fill.background()

    add_text(sl, title, 0.35, 0.1, 10, 0.7, size=28, bold=True,
             color=ACCENT, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(sl, subtitle, 0.35, 0.62, 11, 0.45, size=14,
                 color=TEXT_SUB, align=PP_ALIGN.LEFT)


def bullet_slide(sl, items, x=0.5, y=1.25, w=12.3, size=16,
                 bullet="▸ ", color=TEXT_MAIN, spacing_pt=2):
    tf_box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5.8))
    tf_box.word_wrap = True
    tf = tf_box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(spacing_pt)
        run = p.add_run()
        run.text = bullet + item
        run.font.size  = Pt(size)
        run.font.color.rgb = color


def two_col(sl, left_items, right_items, left_color=TEXT_MAIN, right_color=GREEN,
            size=15, y=1.3):
    bullet_slide(sl, left_items,  x=0.4,  y=y, w=5.9, size=size, color=left_color)
    bullet_slide(sl, right_items, x=6.9,  y=y, w=5.9, size=size, color=right_color)


def divider(sl, y, color=ACCENT, width_pt=0.75):
    line = sl.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Pt(width_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def code_box(sl, code_lines, x, y, w, h, font_size=11):
    box = add_box(sl, x, y, w, h, fill=RGBColor(0x11, 0x11, 0x1B),
                  line_color=RGBColor(0x45, 0x47, 0x5A))
    txb = sl.shapes.add_textbox(Inches(x+0.15), Inches(y+0.1),
                                 Inches(w-0.3), Inches(h-0.2))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = GREEN
        run.font.name = "Courier New"


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()

# gradient-ish accent bar on left
bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text(sl, "nf-PipelineBot", 0.55, 1.6, 12, 1.4,
         size=52, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
add_text(sl, "A Nextflow Wrapper for PipelineBot", 0.55, 3.1, 11, 0.7,
         size=24, color=TEXT_MAIN, align=PP_ALIGN.LEFT)
add_text(sl,
         "Reproducible bioinformatics pipelines with uv, SLURM, GPU acceleration,\n"
         "and Nextflow's built-in reporting — no changes to PipelineBot required.",
         0.55, 3.85, 11, 1.0, size=16, color=TEXT_SUB, align=PP_ALIGN.LEFT)

add_text(sl, "github.com/huimingd/nf-PipelineBot",
         0.55, 6.6, 10, 0.55, size=13, color=TEXT_SUB,
         italic=True, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Overview
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Overview",
             "What is nf-PipelineBot and why use it?")

add_box(sl, 0.4, 1.3, 12.4, 1.6, fill=RGBColor(0x28, 0x29, 0x3E),
        line_color=MAUVE)
add_text(sl,
         "nf-PipelineBot wraps PipelineBot — a Python task-execution framework with "
         "resource monitoring — inside a Nextflow workflow.  It clones PipelineBot at "
         "runtime, installs its locked uv environment, and runs it, giving you "
         "Nextflow's reproducibility and SLURM integration with zero changes to "
         "PipelineBot itself.",
         0.7, 1.4, 11.8, 1.4, size=15, color=TEXT_MAIN)

divider(sl, 3.1)

cards = [
    ("Auto-clone",     ACCENT,  "Clones PipelineBot at any branch, tag, or commit"),
    ("Locked env",     GREEN,   "uv sync --locked guarantees reproducibility"),
    ("SLURM support",  YELLOW,  "Forwards all SLURM flags to PipelineBot"),
    ("GPU / Parabricks", RED,   "pbrun fq2bam in Singularity on GPU nodes"),
]
for i, (title, col, desc) in enumerate(cards):
    cx = 0.4 + i * 3.22
    add_box(sl, cx, 3.25, 3.05, 1.9, line_color=col)
    add_text(sl, title, cx+0.15, 3.35, 2.75, 0.5, size=14, bold=True, color=col)
    add_text(sl, desc,  cx+0.15, 3.85, 2.75, 1.1, size=12, color=TEXT_MAIN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Features
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Features")

features_left = [
    "Automatic Repository Cloning at runtime",
    "Reproducible uv --locked environment",
    "Custom YAML config or built-in default",
    "SLURM pass-through (all flags forwarded)",
]
features_right = [
    "GPU acceleration via Parabricks / Singularity",
    "Publishes results + pipeline_info/ reports",
    "Modular design under modules/ & workflows/",
    "Wave + Conda integration in nextflow.config",
]
two_col(sl, features_left, features_right, left_color=ACCENT, right_color=GREEN,
        size=16, y=1.3)

divider(sl, 5.5)
add_text(sl, "All features are configurable via Nextflow --params without modifying any source file.",
         0.4, 5.6, 12.5, 0.6, size=13, color=TEXT_SUB, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Requirements & Quick Start
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Requirements & Quick Start")

add_text(sl, "Requirements", 0.4, 1.25, 5.5, 0.45, size=17, bold=True, color=ACCENT)
reqs = ["Nextflow ≥ 23.04", "uv on $PATH (or via conda)", "Git"]
bullet_slide(sl, reqs, x=0.4, y=1.7, w=5.5, size=15, color=TEXT_MAIN)

add_text(sl, "Common invocations", 6.5, 1.25, 6.5, 0.45, size=17, bold=True, color=ACCENT)
code_box(sl, [
    "# Default config",
    "nextflow run main.nf",
    "",
    "# Custom YAML",
    "nextflow run main.nf \\",
    "  --config_file /path/config.yaml",
    "",
    "# Specific revision",
    "nextflow run main.nf \\",
    "  --pipelinebot_revision develop",
], x=6.5, y=1.7, w=6.5, h=3.5, font_size=12)

add_text(sl, "SLURM execution", 0.4, 4.1, 5.5, 0.45, size=17, bold=True, color=YELLOW)
code_box(sl, [
    "nextflow run main.nf \\",
    "  --use_slurm true \\",
    "  --slurm_partition compute \\",
    "  --slurm_account myproject \\",
    "  --slurm_time_limit '08:00:00' \\",
    "  --slurm_mem_gb 32",
], x=0.4, y=4.55, w=5.8, h=2.6, font_size=12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Parameters (core + data)
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Parameters — Core & Data")

headers = ["Parameter", "Default", "Description"]
rows = [
    ("pipelinebot_repo",      "…/PipelineBot.git",       "PipelineBot repository URL"),
    ("pipelinebot_revision",  "main",                     "Branch, tag, or commit"),
    ("config_file",           "null",                     "Path to YAML config (uses built-in if unset)"),
    ("output_dir",            "results",                  "Directory for published outputs"),
    ("log_level",             "INFO",                     "Logging level passed to main.py"),
    ("reference_genome",      "null",                     "Reference genome path (passed to YAML)"),
    ("annotation_file",       "null",                     "Annotation GTF path (passed to YAML)"),
    ("fastq_files",           "[]",                       "List of input FASTQ files"),
    ("bam_files",             "[]",                       "List of input BAM files"),
]
col_w = [3.1, 2.4, 6.5]
col_x = [0.35, 3.5, 5.95]
row_h = 0.48
y0 = 1.25

# header row
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_box(sl, cx, y0, cw, row_h, fill=RGBColor(0x1E, 0x3A, 0x5F),
            line_color=ACCENT, line_width_pt=0.5)
    add_text(sl, hdr, cx+0.08, y0+0.06, cw-0.1, row_h-0.1,
             size=13, bold=True, color=ACCENT)

for ri, row in enumerate(rows):
    yy = y0 + (ri+1)*row_h
    bg = RGBColor(0x28, 0x29, 0x3E) if ri % 2 == 0 else BOX_BG
    for ci, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_box(sl, cx, yy, cw, row_h, fill=bg,
                line_color=RGBColor(0x45, 0x47, 0x5A), line_width_pt=0.3)
        color = YELLOW if ci == 0 else (GREEN if ci == 1 else TEXT_MAIN)
        add_text(sl, val, cx+0.08, yy+0.06, cw-0.1, row_h-0.1,
                 size=11, color=color)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Parameters (SLURM)
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Parameters — SLURM Pass-through")

slurm_rows = [
    ("use_slurm",             "false",               "Enable SLURM execution inside PipelineBot"),
    ("slurm_partition",       "null",                "SLURM partition / queue"),
    ("slurm_account",         "null",                "SLURM project account"),
    ("slurm_time_limit",      "04:00:00",            "Wall-clock limit per job"),
    ("slurm_mem_gb",          "null",                "Memory per job (GB)"),
    ("slurm_nodes",           "null",                "Nodes per job"),
    ("slurm_ntasks_per_node", "null",                "MPI tasks per node"),
    ("slurm_work_dir",        "/tmp/pipeline_slurm", "Shared dir for job scripts & logs"),
    ("slurm_python",          "python",              "Python binary on compute nodes"),
    ("slurm_poll_interval",   "10",                  "Seconds between squeue polls"),
]
y0 = 1.25
for ri, row in enumerate(slurm_rows):
    yy = y0 + (ri+1)*row_h if ri > 0 else y0 + row_h
    if ri == 0:
        # header
        for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
            add_box(sl, cx, y0, cw, row_h, fill=RGBColor(0x1E, 0x3A, 0x5F),
                    line_color=ACCENT, line_width_pt=0.5)
            add_text(sl, hdr, cx+0.08, y0+0.06, cw-0.1, row_h-0.1,
                     size=13, bold=True, color=ACCENT)

    yy = y0 + (ri+1)*row_h
    bg = RGBColor(0x28, 0x29, 0x3E) if ri % 2 == 0 else BOX_BG
    for ci, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_box(sl, cx, yy, cw, row_h, fill=bg,
                line_color=RGBColor(0x45, 0x47, 0x5A), line_width_pt=0.3)
        color = YELLOW if ci == 0 else (GREEN if ci == 1 else TEXT_MAIN)
        add_text(sl, val, cx+0.08, yy+0.06, cw-0.1, row_h-0.1,
                 size=11, color=color)

add_text(sl,
         "All slurm_* params map directly to PipelineBot's --slurm-* CLI flags.",
         0.4, 6.8, 12, 0.45, size=13, color=TEXT_SUB, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Pipeline Config YAML
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Pipeline Config YAML",
             "Pass --config_file to override the built-in bioinformatics example")

code_box(sl, [
    "tasks:",
    "  - alignment",
    "  - variant_calling",
    "  - rna_seq",
    "",
    "executor:",
    "  reference_genome: reference_genome.fa",
    "  annotation_file:  annotations.gtf",
    "  output_dir:       bioinformatics_results",
    "",
    "alignment:",
    "  resource_config: { cpus: 4, memory_gb: 8.0, max_processes: 2 }",
    "  fastq_files: [sample1.fastq, sample2.fastq]",
    "  aligner: bwa   # bwa | bowtie2 | star | minimap2 | parabricks",
], x=0.35, y=1.25, w=6.2, h=5.8, font_size=11)

code_box(sl, [
    "variant_calling:",
    "  resource_config: { cpus: 8, memory_gb: 16.0 }",
    "  caller: gatk   # gatk | freebayes | samtools",
    "  bam_files: []  # auto-populated from alignment",
    "",
    "rna_seq:",
    "  resource_config: { cpus: 6, memory_gb: 12.0 }",
    "  quantification_method: salmon",
    "  # options: salmon | kallisto | featurecounts",
    "",
    "# Per-task SLURM override (optional)",
    "# alignment:",
    "#   slurm_config:",
    "#     partition: compute",
    "#     mem_gb: 8.0",
    "#     time_limit: '02:00:00'",
], x=6.75, y=1.25, w=6.2, h=5.8, font_size=11)

add_text(sl,
         "Tasks run in the order listed under tasks:. "
         "Remove entries to skip steps. BAM files flow automatically from alignment → variant_calling.",
         0.35, 6.9, 12.5, 0.4, size=12, color=TEXT_SUB, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — GPU Acceleration (Parabricks)
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "GPU-Accelerated Alignment with Parabricks",
             "Set  aligner: parabricks  to run pbrun fq2bam inside Singularity")

bullets = [
    "Container launched with  singularity exec --nv  (GPU pass-through)",
    "scratch_dir  on the host is bind-mounted to  /scratch  inside the container",
    "All input/output paths should be expressed relative to  /scratch",
    "Per-task slurm_config block selects the GPU partition and device count",
]
bullet_slide(sl, bullets, x=0.4, y=1.3, w=12.4, size=15, color=TEXT_MAIN)

code_box(sl, [
    "executor:",
    "  reference_genome: /scratch/Genomes/hg38.fa",
    "  output_dir:       /RUN/alignment",
    "",
    "alignment:",
    "  aligner:        parabricks",
    "  fastq_files:    [/scratch/sample_R1.fastq]",
    "  fastq_r2_files: [/scratch/sample_R2.fastq]",
    "  container_path: /host/path/clara-parabricks.sif",
    "  scratch_dir:    /host/path/to/scratch",
    "  slurm_config:",
    "    partition: gpu_partition",
    "    mem_gb:    64.0",
    "    time_limit: '06:00:00'",
    "    gpu:       h200:1   # or l40s:1, h100:2 …",
    "    modules:   ['cuda/13.0.1', 'apptainer/1.1.9']",
    "    work_dir:  /host/path/slurm_jobs",
], x=0.4, y=3.2, w=8.5, h=4.0, font_size=11)

add_text(sl, "Run with:", 9.3, 3.2, 3.7, 0.4, size=13, bold=True, color=ACCENT)
code_box(sl, [
    "nextflow run main.nf \\",
    "  --config_file \\",
    "  /path/alignment_gpu.yaml",
], x=9.3, y=3.65, w=3.65, h=1.2, font_size=12)

add_text(sl, "Supported GPU types", 9.3, 5.05, 3.7, 0.4, size=13, bold=True, color=ACCENT)
bullet_slide(sl, ["h200:1", "h100:2", "l40s:1"],
             x=9.3, y=5.45, w=3.65, size=13, color=GREEN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Workflow Structure
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Workflow Structure")

# process boxes
steps = [
    ("DOWNLOAD_PIPELINEBOT",  "git clone + checkout",  ACCENT),
    ("SETUP_UV_ENVIRONMENT",  "uv sync --locked",       GREEN),
    ("RUN_PIPELINEBOT",       "uv run python main.py",  MAUVE),
]
bx, by, bw, bh = 1.2, 1.35, 4.2, 1.0
for i, (name, detail, col) in enumerate(steps):
    yy = by + i * 1.6
    add_box(sl, bx, yy, bw, bh, line_color=col)
    add_text(sl, name,   bx+0.18, yy+0.08, bw-0.3, 0.5, size=15, bold=True, color=col)
    add_text(sl, detail, bx+0.18, yy+0.52, bw-0.3, 0.4, size=12, color=TEXT_SUB)
    if i < len(steps)-1:
        arrow = sl.shapes.add_shape(1, Inches(bx+bw/2-0.05), Inches(yy+bh),
                                     Inches(0.1), Inches(0.6))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEXT_SUB
        arrow.line.fill.background()

# output tree
add_text(sl, "results/", 6.5, 1.35, 6.5, 0.45, size=16, bold=True, color=YELLOW)
tree = [
    ("pipeline_results.json", "Pipeline task outcomes (JSON)", GREEN),
    ("pipeline.log",          "Full execution log",            GREEN),
    ("pipeline_info/",        "Nextflow reporting outputs",    ACCENT),
    ("  execution_timeline.html", "Per-process wall-clock",    TEXT_SUB),
    ("  execution_report.html",   "CPU / memory usage",        TEXT_SUB),
    ("  execution_trace.txt",     "Tab-separated task trace",  TEXT_SUB),
    ("  pipeline_dag.svg",        "Workflow DAG diagram",      TEXT_SUB),
]
for i, (fname, desc, col) in enumerate(tree):
    yy = 1.9 + i * 0.65
    prefix = "├── " if i < len(tree)-1 and not fname.startswith("  ") else \
             ("└── " if not fname.startswith("  ") else "│   ")
    add_text(sl, prefix + fname.strip(), 6.5, yy, 3.5, 0.55,
             size=13, color=col)
    add_text(sl, desc, 10.0, yy, 3.2, 0.55, size=12, color=TEXT_SUB, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Modules
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Modules & Workflow Layout")

module_info = [
    ("modules/download_pipelinebot.nf",
     "Clones the PipelineBot repo at the specified revision (git clone + checkout).",
     ACCENT),
    ("modules/setup_uv_environment.nf",
     "Runs uv sync --locked and verifies psutil/yaml are importable.",
     GREEN),
    ("modules/run_pipelinebot.nf",
     "Executes uv run python main.py with config, output paths, log level, and all SLURM args.\n"
     "Writes outputs to Nextflow's WORK_DIR so publishDir can copy them to results/.",
     MAUVE),
    ("workflows/pipelinebot.nf",
     "Orchestrates the three modules in sequence and emits results + log_file channels.",
     YELLOW),
]
for i, (path, desc, col) in enumerate(module_info):
    yy = 1.3 + i * 1.45
    add_box(sl, 0.35, yy, 12.6, 1.3, line_color=col)
    add_text(sl, path, 0.55, yy+0.1, 7.5, 0.5, size=14, bold=True, color=col)
    add_text(sl, desc, 0.55, yy+0.58, 12.0, 0.65, size=12, color=TEXT_MAIN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Reporting & Wave
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Reporting & Wave + Conda")

add_text(sl, "Nextflow Execution Reports", 0.4, 1.25, 6.2, 0.45,
         size=17, bold=True, color=ACCENT)
report_rows = [
    ("execution_timeline.html", "Per-process wall-clock timeline"),
    ("execution_report.html",   "Resource usage — CPU, memory, I/O"),
    ("execution_trace.txt",     "Tab-separated trace of every task"),
    ("pipeline_dag.svg",        "Directed acyclic graph of the workflow"),
]
for i, (fname, desc) in enumerate(report_rows):
    yy = 1.75 + i * 0.7
    add_box(sl, 0.4, yy, 6.2, 0.62, line_color=RGBColor(0x45,0x47,0x5A))
    add_text(sl, fname, 0.58, yy+0.06, 2.8, 0.48, size=12, color=GREEN)
    add_text(sl, desc,  3.4,  yy+0.06, 3.1, 0.48, size=12, color=TEXT_MAIN)
add_text(sl, "Written to  results/pipeline_info/  after every run.",
         0.4, 4.7, 6.2, 0.4, size=12, color=TEXT_SUB, italic=True)

add_text(sl, "Wave + Conda Integration", 7.1, 1.25, 6.0, 0.45,
         size=17, bold=True, color=MAUVE)
add_box(sl, 7.1, 1.75, 6.0, 3.5, line_color=MAUVE)
wave_text = (
    "nextflow.config enables Wave with strategy = ['conda'].\n\n"
    "When wave.enabled = true, Nextflow can automatically build "
    "and cache OCI container images from Conda environment "
    "definitions on processes that declare a conda directive.\n\n"
    "No manual container management required — Wave handles "
    "image building and caching transparently."
)
add_text(sl, wave_text, 7.3, 1.85, 5.7, 3.2, size=13, color=TEXT_MAIN)

code_box(sl, [
    "# nextflow.config",
    "wave {",
    "  enabled  = true",
    "  strategy = ['conda']",
    "}",
    "conda {",
    "  enabled  = true",
    "  channels = ['conda-forge',",
    "              'bioconda']",
    "}",
], x=7.1, y=5.3, w=6.0, h=1.95, font_size=11)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Entry Points
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()
slide_header(sl, "Entry Points")

add_box(sl, 0.4, 1.3, 5.9, 3.0, line_color=ACCENT)
add_text(sl, "main.nf  (recommended)", 0.6, 1.4, 5.5, 0.5,
         size=16, bold=True, color=ACCENT)
add_text(sl,
         "Delegates to workflows/pipelinebot.nf which includes all three modules.\n\n"
         "Split structure makes individual steps easier to reuse, override, or extend "
         "in downstream workflows via  include { … } from './workflows/pipelinebot'.",
         0.6, 1.95, 5.5, 2.1, size=13, color=TEXT_MAIN)

add_box(sl, 7.0, 1.3, 5.9, 3.0, line_color=GREEN)
add_text(sl, "main_single.nf", 7.2, 1.4, 5.5, 0.5,
         size=16, bold=True, color=GREEN)
add_text(sl,
         "Portable single-file workflow — all three processes and the workflow block "
         "live in one file with no external module dependencies.\n\n"
         "Useful for sharing a self-contained script or running without the full "
         "repository checkout.",
         7.2, 1.95, 5.5, 2.1, size=13, color=TEXT_MAIN)

divider(sl, 4.55)

add_text(sl,
         "Both entry points are functionally identical and share the same nextflow.config — "
         "choose based on whether you need modularity or portability.",
         0.4, 4.65, 12.5, 0.7, size=14, color=TEXT_SUB, italic=True,
         align=PP_ALIGN.CENTER)

code_box(sl, [
    "nextflow run main.nf        # modular (recommended)",
    "nextflow run main_single.nf # portable single-file",
], x=2.5, y=5.45, w=8.3, h=1.2, font_size=13)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Summary
# ─────────────────────────────────────────────────────────────────────────────
sl = new_slide()

bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()

add_text(sl, "Summary", 0.55, 1.2, 11, 0.9,
         size=40, bold=True, color=GREEN, align=PP_ALIGN.LEFT)

summary = [
    "Clone PipelineBot at any revision — fully reproducible",
    "uv --locked environment ensures dependency consistency",
    "Custom YAML config drives all pipeline logic",
    "SLURM params forwarded transparently to PipelineBot",
    "GPU jobs via Parabricks + Singularity on GPU partitions",
    "pipeline_info/ reports (timeline, trace, DAG) auto-generated",
    "Wave + Conda for zero-touch container provisioning",
    "Two entry points: modular main.nf or portable main_single.nf",
]
bullet_slide(sl, summary, x=0.55, y=2.2, w=12.3, size=15,
             bullet="✓  ", color=TEXT_MAIN, spacing_pt=4)

add_text(sl, "github.com/huimingd/nf-PipelineBot",
         0.55, 6.7, 12, 0.5, size=13, color=TEXT_SUB,
         italic=True, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
out = "/orcd/data/ki/001/core/bmc/users/huiming/sourcecode/nf-PipelineBot/nf-PipelineBot.pptx"
prs.save(out)
print(f"Saved: {out}")
